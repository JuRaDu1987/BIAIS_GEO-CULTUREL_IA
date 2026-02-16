from pptx import Presentation
import os

pptx_path = r"c:\Users\stagiaire\Desktop\PROJET_IA\PROJET2\Biais IA Géo-Culturel.pptx"

if not os.path.exists(pptx_path):
    print(f"Error: File not found at {pptx_path}")
    exit(1)

prs = Presentation(pptx_path)

print(f"Presentation: {os.path.basename(pptx_path)}")
print(f"Number of slides: {len(prs.slides)}\n")

for i, slide in enumerate(prs.slides):
    print(f"--- Slide {i+1} ---")
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text = shape.text.strip()
            if text:
                print(text)
    print("\n")
