from pptx import Presentation
import os
import sys

# Set encoding to utf-8 for output
sys.stdout.reconfigure(encoding='utf-8')

pptx_path = r"c:\Users\stagiaire\Desktop\PROJET_IA\PROJET2\Biais IA Géo-Culturel.pptx"

if not os.path.exists(pptx_path):
    print(f"Error: File not found at {pptx_path}")
    exit(1)

prs = Presentation(pptx_path)

for i, slide in enumerate(prs.slides):
    if i + 1 in [7, 8]:
        print(f"--- Slide {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    print(text.encode('ascii', 'ignore').decode('ascii'))
        print("\n")
